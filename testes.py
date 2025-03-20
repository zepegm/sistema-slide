from pptx import Presentation

def extract_text_with_formatting(paragraph):
    html_text = ""
    
    for run in paragraph.runs:
        text = run.text
        if not text.strip():
            continue  # Ignorar textos vazios

        style = ""

        if run.font.bold:
            text = f"<b>{text}</b>"
        if run.font.italic:
            text = f"<i>{text}</i>"
        if run.font.underline:
            text = f"<u>{text}</u>"

        try:
            if run.font.color and run.font.color.rgb:
                color = run.font.color.rgb
                style += f"color: #{color};"
        except:
            pass

        if style:
            text = f'<span style="{style}">{text}</span>'

        html_text += text

    return html_text

def pptx_to_html(pptx_file, output_html):
    prs = Presentation(pptx_file)
    html_content = """
    <html>
    <head>
        <style>
            body { font-family: Arial, sans-serif; }
            .slide { margin-bottom: 40px; padding: 10px; border-bottom: 2px solid #ccc; }
            h2 { color: #333; }
        </style>
    </head>
    <body>
    """

    for i, slide in enumerate(prs.slides):
        html_content += f'<div class="slide"><h2>Slide {i+1}</h2>'

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    formatted_text = extract_text_with_formatting(paragraph)
                    if formatted_text:
                        html_content += f"<p>{formatted_text}</p>"

        html_content += "</div>"

    html_content += "</body></html>"

    with open(output_html, "w", encoding="utf-8") as f:
        f.write(html_content)

    print(f"✅ Conversão concluída! Arquivo salvo como {output_html}")

# Exemplo de uso:
pptx_to_html(r"C:\Users\giuseppe.manzella\OneDrive - Secretaria da Educação do Estado de São Paulo\IGREJA\Músicas\Escuro\Noiva Adornada e Pronta.pptx", "saida.html")