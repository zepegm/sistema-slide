from pptx import Presentation

dir = r'C:\Users\giuseppe.manzella\OneDrive - Secretaria da Educação do Estado de São Paulo\IGREJA\Músicas\Escuro\Autor da Minha Fé (UMADECAP).pptx'

file = open(f'{dir}', 'rb')
prs = Presentation(file)

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:

                print(run.text)
