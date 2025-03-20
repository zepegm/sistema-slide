from pptx import Presentation
import win32com.client
import pythoncom
import os
import re

def substituir_quebra_linha(texto):
    # Expressão regular para encontrar \u000b
    padrao = re.compile(r'\u000b([A-Z]?)')
    # Substituir \u000b por <br> se seguido por uma letra maiúscula, ou por um espaço caso contrário
    texto_modificado = padrao.sub(lambda m: '<br>' + m.group(1) if m.group(1) else ' ', texto)
    return texto_modificado

def getListText(dir):
    file = open(f'{dir}', 'rb')

    prs = Presentation(file)


    text_runs = []

    slide_pos = 0

    for slide in prs.slides:

        anotacao = ''
        if slide.has_notes_slide:
            anotacao = slide.notes_slide.notes_text_frame.text

        slide_pos += 1

        if slide_pos == 1:
            continue

        if slide_pos == len(prs.slides):
            continue

        text_slide = ""
        plain_text = ""

        key_b = False
        key_u = False
        key_m = False


        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            cont = 0
            
            for paragraph in shape.text_frame.paragraphs:

                total_runs = len(paragraph.runs)

                for run in paragraph.runs:

                    try:
                        if not key_m and str(run.font.color.theme_color) == 'ACCENT_4 (8)':
                            text_slide +='<mark class="cdx-marker">'
                            key_m = True
                        elif key_m and str(run.font.color.theme_color) != 'ACCENT_4 (8)': 
                            text_slide += "</mark>"
                            key_m = False
                    except:
                        if key_m:
                            text_slide += "</mark>"
                            key_m = False

                    if (not key_b and run.font.bold and not key_m):
                        text_slide +="<b>"
                        key_b = True
                    elif (key_b and not run.font.bold):
                        text_slide += "</b>"
                        key_b = False

                    if (not key_u and run.font.underline and not key_m):
                        text_slide +='<u class="cdx-underline">'
                        key_u = True
                    elif (key_u and not run.font.underline):
                        text_slide += "</u>"
                        key_u = False

                    if cont > 0:
                        text_slide += "<br>" + run.text.strip()
                        plain_text += " " + run.text.strip()
                    else:
                        text_slide += run.text.strip()
                        plain_text += run.text.strip()
                        cont += 1

                if (text_slide != '<b>' and run.text != paragraph.runs[total_runs -1].text):
                    text_slide += "<br>"

                plain_text += "<br>"


        if key_b:
            text_slide += "</b>"
        if key_u:
            text_slide += "</u>"

        text_slide = text_slide.replace('  ', ' ').replace(" </mark> ", "</mark>")
        plain_text = plain_text.replace('  ', ' ')

        print(text_slide)

        text_runs.append({'pos':slide_pos - 1, 'text-slide':text_slide, 'subtitle':plain_text, 'anotacao':anotacao})


    # antes de retornar o texto extrair a capa do slide e converter em jpg
    
    path_img = os.path.dirname(os.path.realpath(__file__)) + '\\static\\images\\SlidesPPTX'
    image_path = os.path.join(path_img, "temp_capa.jpg")
    
    try:
        Application = win32com.client.Dispatch("PowerPoint.Application", pythoncom.CoInitialize())
        # Open the presentation without making it visible
        pptx = Application.Presentations.Open(FileName=dir, WithWindow=False)    
    
        pptx.Slides[0].Export(image_path, "JPG")

    except Exception as e:
        print(f"An error occurred: {e}")

    file.close()
    return text_runs


def getListTextHarpa(dir):

    file = open(f'{dir}', 'rb')

    prs = Presentation(file)

    text_runs = []

    slide_pos = 0

    for slide in prs.slides:

        anotacao = ''
        if slide.has_notes_slide:
            anotacao = slide.notes_slide.notes_text_frame.text

        slide_pos += 1

        if slide_pos == 1:
            continue

        text_slide = "<b>"
        plain_text = ""
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
  
            for paragraph in shape.text_frame.paragraphs:

                try:
                    if str(paragraph.runs[0].font.color.rgb) == '0070C0':
                        text_slide += '<span class="cdx-num">' + paragraph.runs[0].text + '</span>' + paragraph.text.replace(paragraph.runs[0].text, ' ', 1).replace('\u000b', '<br>').rstrip()
                    elif str(paragraph.runs[0].font.color.rgb) == 'FF0000':
                        if paragraph.text[0:1].isdigit():
                            text_slide += '<span class="red">' + paragraph.runs[0].text + '</span>' + paragraph.text.replace(paragraph.runs[0].text, ' ', 1).replace('\u000b', '<br>').rstrip()
                        else:
                            text_slide += '<span class="red">' + paragraph.text.replace('\u000b', '<br>').lstrip() + '</span>'
                    else:
                        text_slide += paragraph.text.replace('\u000b', '<br>').lstrip()
                except:
                    text_slide += paragraph.text.replace('\u000b', '<br>').lstrip()

                plain_text += substituir_quebra_linha(paragraph.text)

                if (text_slide != '<b>'):
                    text_slide += "<br>"
                    plain_text += "<br>"

        text_slide += '</b>'

        text_runs.append({'pos':slide_pos - 1, 'text-slide':text_slide, 'subtitle':plain_text, 'anotacao':anotacao})
    
    file.close()
    return text_runs